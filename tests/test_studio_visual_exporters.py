from __future__ import annotations

from pathlib import Path

import fitz
import pytest
from PIL import Image, ImageStat
from pptx import Presentation

from open_notebook.studio.exporters import export_infographic, export_slide_deck
from open_notebook.studio.schemas import (
    InfographicDocument,
    SlideDeckDocument,
    parse_artifact_document,
)


def _slide_deck() -> SlideDeckDocument:
    document = parse_artifact_document(
        "slide_deck",
        {
            "artifact_type": "slide_deck",
            "title": "Local Evidence Studio",
            "audience": "Private researchers",
            "slides": [
                {
                    "title": "Grounded generation",
                    "bullets": [
                        "Artifacts remain tied to selected sources.",
                        "Local models can produce validated documents.",
                    ],
                    "speaker_notes": "Explain why source ownership matters.",
                    "visual_direction": "Use a simple evidence flow.",
                    "citations": ["[S1]", "[S2]"],
                },
                {
                    "title": "Private by default",
                    "bullets": ["No hosted rendering service is required."],
                    "citations": ["[S2]"],
                },
            ],
        },
    )
    assert isinstance(document, SlideDeckDocument)
    return document


def _infographic(orientation: str = "portrait") -> InfographicDocument:
    document = parse_artifact_document(
        "infographic",
        {
            "artifact_type": "infographic",
            "title": "Evidence at a glance",
            "orientation": orientation,
            "panels": [
                {
                    "kind": "metric",
                    "heading": "Source coverage",
                    "value": "95%",
                    "body": "Resolved citation target",
                    "citations": ["[S1]"],
                },
                {
                    "kind": "process",
                    "heading": "Workflow",
                    "body": "Collect, validate, render, and review.",
                    "citations": ["[S2]"],
                },
                {
                    "kind": "comparison",
                    "heading": "Ownership",
                    "body": "Local files and models remain under owner control.",
                },
            ],
        },
    )
    assert isinstance(document, InfographicDocument)
    return document


def _dense_infographic(orientation: str) -> InfographicDocument:
    document = parse_artifact_document(
        "infographic",
        {
            "artifact_type": "infographic",
            "title": "Twenty grounded findings",
            "orientation": orientation,
            "panels": [
                {
                    "kind": "metric" if index % 3 == 0 else "text",
                    "heading": f"Finding {index + 1}",
                    "value": f"{index + 1}%" if index % 3 == 0 else "",
                    "body": "A compact source-grounded finding.",
                    "citations": [f"[S{(index % 4) + 1}]"],
                }
                for index in range(20)
            ],
        },
    )
    assert isinstance(document, InfographicDocument)
    return document


def _assert_nonblank_image(path: Path, expected_size: tuple[int, int]) -> None:
    with Image.open(path) as image:
        assert image.size == expected_size
        assert image.mode in {"RGB", "RGBA"}
        extrema = ImageStat.Stat(image.convert("RGB")).extrema
        assert any(low != high for low, high in extrema)


def test_slide_deck_exports_editable_pptx_with_notes_and_citations(tmp_path):
    pptx_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"

    export_slide_deck(_slide_deck(), pptx_path, pdf_path)

    presentation = Presentation(pptx_path)
    assert presentation.slide_width == 12192000
    assert presentation.slide_height == 6858000
    assert len(presentation.slides) == 3
    assert "Local Evidence Studio" in " ".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    content_slide = presentation.slides[1]
    slide_text = " ".join(
        shape.text for shape in content_slide.shapes if hasattr(shape, "text")
    )
    assert "Grounded generation" in slide_text
    assert "Artifacts remain tied to selected sources." in slide_text
    assert "[S1] [S2]" in slide_text
    notes = content_slide.notes_slide.notes_text_frame.text
    assert "Explain why source ownership matters." in notes
    assert "Use a simple evidence flow." in notes
    assert "[S1] [S2]" in notes

    pdf = fitz.open(pdf_path)
    assert pdf.page_count == 3
    pixmap = pdf[1].get_pixmap(matrix=fitz.Matrix(0.25, 0.25), alpha=False)
    assert len(set(pixmap.samples)) > 8
    pdf.close()


@pytest.mark.parametrize(
    ("orientation", "expected_size"),
    [
        ("portrait", (1200, 1800)),
        ("landscape", (1800, 1200)),
        ("square", (1400, 1400)),
    ],
)
def test_infographic_exports_nonblank_png_and_pdf(
    tmp_path,
    orientation,
    expected_size,
):
    png_path = tmp_path / f"infographic-{orientation}.png"
    pdf_path = tmp_path / f"infographic-{orientation}.pdf"

    export_infographic(_infographic(orientation), png_path, pdf_path)

    _assert_nonblank_image(png_path, expected_size)
    pdf = fitz.open(pdf_path)
    assert pdf.page_count == 1
    assert pdf[0].rect.width > 0
    assert pdf[0].rect.height > 0
    pdf.close()


@pytest.mark.parametrize(
    ("orientation", "expected_size"),
    [
        ("portrait", (1200, 1800)),
        ("landscape", (1800, 1200)),
        ("square", (1400, 1400)),
    ],
)
def test_infographic_exports_schema_maximum_panel_count(
    tmp_path,
    orientation,
    expected_size,
):
    png_path = tmp_path / f"dense-{orientation}.png"
    pdf_path = tmp_path / f"dense-{orientation}.pdf"

    export_infographic(_dense_infographic(orientation), png_path, pdf_path)

    _assert_nonblank_image(png_path, expected_size)
    with fitz.open(pdf_path) as pdf:
        assert pdf.page_count == 1


def test_visual_exporters_reject_the_wrong_document(tmp_path):
    with pytest.raises(TypeError, match="SlideDeckDocument"):
        export_slide_deck(
            _infographic(),
            tmp_path / "wrong.pptx",
            tmp_path / "wrong.pdf",
        )

    with pytest.raises(TypeError, match="InfographicDocument"):
        export_infographic(
            _slide_deck(),
            tmp_path / "wrong.png",
            tmp_path / "wrong.pdf",
        )
