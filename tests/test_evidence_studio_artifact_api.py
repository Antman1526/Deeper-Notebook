"""Evidence Studio artifact API tests."""

from __future__ import annotations

import json
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, MagicMock

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation

from api.routers import studio as studio_mod
from deeper_notebook.exceptions import NotFoundError
from deeper_notebook.studio.generation import service as artifact_generation_service


def _client() -> TestClient:
    app = FastAPI()
    app.include_router(studio_mod.router, prefix="/api")
    return TestClient(app)


@pytest.fixture(autouse=True)
def _isolate_background_evaluation(monkeypatch):
    """Keep router tests focused on API behavior, not live evaluation persistence."""
    for name in (
        "DEEPER_NOTEBOOK_EVIDENCE_STUDIO",
        "DN_EVIDENCE_STUDIO",
        "DEEPER_NOTEBOOK_EVIDENCE_STUDIO",
    ):
        monkeypatch.delenv(name, raising=False)
    monkeypatch.setattr(
        artifact_generation_service,
        "_schedule_artifact_evaluation",
        lambda **_kwargs: None,
    )


def _minimal_document(artifact_type: str, title: str | None = None) -> dict:
    document_title = title or artifact_type.replace("_", " ").title()
    if artifact_type in {"report", "study_guide", "briefing", "faq", "timeline"}:
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "sections": [{"heading": "Evidence", "body": "Grounded answer."}],
        }
    if artifact_type in {"course_pack", "training_guide"}:
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "audience": "Learners",
            "learning_outcomes": ["Apply the source evidence"],
            "modules": [
                {
                    "title": "Foundation",
                    "lessons": [{"title": "Evidence", "content": "Grounded lesson."}],
                }
            ],
        }
    if artifact_type == "flashcards":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "cards": [{"front": "Question", "back": "Answer"}],
        }
    if artifact_type == "quiz":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "questions": [
                {
                    "prompt": "Question?",
                    "options": [
                        {"id": "A", "text": "Answer"},
                        {"id": "B", "text": "Distractor"},
                    ],
                    "correct_option_id": "A",
                }
            ],
        }
    if artifact_type == "data_table":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "columns": ["Topic", "Evidence"],
            "rows": [{"values": {"Topic": "Studio", "Evidence": "Grounded"}}],
        }
    if artifact_type == "mind_map":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "root": {"label": "Open Notebook Plus", "children": []},
        }
    if artifact_type == "slide_deck":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "slides": [{"title": "Evidence", "bullets": ["Grounded claim"]}],
        }
    if artifact_type == "infographic":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "panels": [{"kind": "text", "heading": "Evidence"}],
        }
    if artifact_type == "podcast_outline":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "cold_open": "Start with evidence.",
            "segments": [{"title": "Grounding", "beats": ["Verify claims"]}],
        }
    if artifact_type == "research_run":
        return {
            "artifact_type": artifact_type,
            "title": document_title,
            "objective": "Investigate the evidence.",
            "stages": [{"title": "Synthesis", "findings": []}],
        }
    raise AssertionError(f"No test document for {artifact_type}")


def _json_chain(
    artifact_type: str,
    *,
    title: str | None = None,
    document: dict | None = None,
) -> MagicMock:
    chain = MagicMock()
    chain.with_structured_output.side_effect = NotImplementedError
    chain.ainvoke = AsyncMock(
        return_value=SimpleNamespace(
            content=json.dumps(document or _minimal_document(artifact_type, title))
        )
    )
    return chain


class _FakeArtifact:
    records: dict[str, "_FakeArtifact"] = {}
    deleted: list[str] = []
    saved_ids: list[str] = []

    def __init__(self, **kwargs):
        self.id = kwargs.pop("id", None)
        self.created = kwargs.pop("created", None)
        self.updated = kwargs.pop("updated", None)
        self.notebook_id = kwargs.pop("notebook_id")
        self.artifact_type = kwargs.pop("artifact_type")
        self.title = kwargs.pop("title")
        self.status = kwargs.pop("status", "pending")
        self.source_ids = kwargs.pop("source_ids", [])
        self.prompt = kwargs.pop("prompt", None)
        self.model_id = kwargs.pop("model_id", None)
        self.provider = kwargs.pop("provider", None)
        self.output_format = kwargs.pop("output_format", None)
        self.output_payload = kwargs.pop("output_payload", {})
        self.citations = kwargs.pop("citations", [])
        self.export_paths = kwargs.pop("export_paths", {})
        self.revision_of_id = kwargs.pop("revision_of_id", None)

    async def save(self):
        if not self.id:
            self.id = f"studio_artifact:{len(self.records) + 1}"
        now = datetime(2026, 6, 23, tzinfo=timezone.utc)
        self.created = self.created or now
        self.updated = now
        self.records[self.id] = self
        self.saved_ids.append(self.id)

    async def delete(self):
        self.deleted.append(self.id)
        self.records.pop(self.id, None)
        return True

    @classmethod
    async def get(cls, artifact_id: str):
        return cls.records[artifact_id]

    @classmethod
    async def get_for_notebook(cls, notebook_id: str):
        return [
            item for item in cls.records.values() if item.notebook_id == notebook_id
        ]

    @classmethod
    async def get_revisions(cls, artifact_id: str):
        return [
            item for item in cls.records.values() if item.revision_of_id == artifact_id
        ]


class _FakeWorkflowRun:
    records: dict[str, "_FakeWorkflowRun"] = {}
    saved_ids: list[str] = []

    def __init__(self, **kwargs):
        self.id = kwargs.pop("id", None)
        self.created = kwargs.pop("created", None)
        self.updated = kwargs.pop("updated", None)
        self.artifact_id = kwargs.pop("artifact_id")
        self.notebook_id = kwargs.pop("notebook_id")
        self.title = kwargs.pop("title")
        self.status = kwargs.pop("status", "queued")
        self.source_ids = kwargs.pop("source_ids", [])
        self.approval_required = kwargs.pop("approval_required", False)
        self.steps = kwargs.pop("steps", [])
        self.command_id = kwargs.pop("command_id", None)

    async def save(self):
        if not self.id:
            self.id = f"studio_workflow_run:{len(self.records) + 1}"
        now = datetime(2026, 6, 23, tzinfo=timezone.utc)
        self.created = self.created or now
        self.updated = now
        self.records[self.id] = self
        self.saved_ids.append(self.id)

    @classmethod
    async def get(cls, run_id: str):
        return cls.records[run_id]

    @classmethod
    async def get_for_artifact(cls, artifact_id: str):
        return [
            item for item in cls.records.values() if item.artifact_id == artifact_id
        ]


def _install_fake_artifacts(monkeypatch):
    _FakeArtifact.records = {}
    _FakeArtifact.deleted = []
    _FakeArtifact.saved_ids = []
    monkeypatch.setattr(studio_mod, "StudioArtifact", _FakeArtifact)
    monkeypatch.setattr(
        studio_mod,
        "repo_query",
        AsyncMock(return_value=[{"id": "notebook:alpha"}]),
    )

    class _NotebookMock:
        @classmethod
        async def get(cls, _notebook_id):
            return cls()

    monkeypatch.setattr(studio_mod, "Notebook", _NotebookMock)


def _install_fake_workflow_runs(monkeypatch):
    _FakeWorkflowRun.records = {}
    _FakeWorkflowRun.saved_ids = []
    monkeypatch.setattr(studio_mod, "StudioWorkflowRun", _FakeWorkflowRun)


def test_artifact_routes_are_hidden_when_feature_flag_disabled(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "0")

    response = _client().get("/api/studio/notebooks/notebook:alpha/artifacts")

    assert response.status_code == 404
    assert response.json()["detail"] == "Evidence Studio is not enabled"


def test_create_artifact_saves_and_returns_response(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)

    response = _client().post(
        "/api/studio/artifacts",
        json={
            "notebook_id": "notebook:alpha",
            "artifact_type": "study_guide",
            "title": "Study guide",
            "source_ids": ["source:one"],
            "prompt": "Make a guide",
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["id"] == "studio_artifact:1"
    assert body["notebook_id"] == "notebook:alpha"
    assert body["artifact_type"] == "study_guide"
    assert body["status"] == "pending"
    assert body["source_ids"] == ["source:one"]


def test_create_training_guide_artifact(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)

    response = _client().post(
        "/api/studio/artifacts",
        json={
            "notebook_id": "notebook:training",
            "artifact_type": "training_guide",
            "title": "User onboarding training",
            "source_ids": ["source:video", "source:pdf"],
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["artifact_type"] == "training_guide"
    assert body["title"] == "User onboarding training"
    assert body["source_ids"] == ["source:video", "source:pdf"]


def test_create_course_pack_artifact(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)

    response = _client().post(
        "/api/studio/artifacts",
        json={
            "notebook_id": "notebook:training",
            "artifact_type": "course_pack",
            "title": "Course Pack",
            "source_ids": ["source:video", "source:pdf", "source:link"],
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["artifact_type"] == "course_pack"
    assert body["title"] == "Course Pack"
    assert body["source_ids"] == ["source:video", "source:pdf", "source:link"]


def test_list_artifacts_for_notebook(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
        ),
        "studio_artifact:2": _FakeArtifact(
            id="studio_artifact:2",
            notebook_id="notebook:beta",
            artifact_type="quiz",
            title="Quiz",
        ),
    }

    response = _client().get("/api/studio/notebooks/notebook:alpha/artifacts")

    assert response.status_code == 200
    body = response.json()
    assert [item["id"] for item in body] == ["studio_artifact:1"]


def test_list_artifacts_for_missing_notebook_returns_404(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    studio_mod.repo_query.return_value = []

    response = _client().get("/api/studio/notebooks/notebook:gone/artifacts")

    assert response.status_code == 404
    assert response.json()["detail"] == "Notebook not found"


def test_list_artifacts_for_notebook_excludes_revision_snapshots(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:primary": _FakeArtifact(
            id="studio_artifact:primary",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
        ),
        "studio_artifact:revision": _FakeArtifact(
            id="studio_artifact:revision",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report revision",
            revision_of_id="studio_artifact:primary",
        ),
    }

    response = _client().get("/api/studio/notebooks/notebook:alpha/artifacts")

    assert response.status_code == 200
    assert [item["id"] for item in response.json()] == ["studio_artifact:primary"]


def test_list_artifact_revisions(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:primary": _FakeArtifact(
            id="studio_artifact:primary",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
        ),
        "studio_artifact:revision": _FakeArtifact(
            id="studio_artifact:revision",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report revision",
            revision_of_id="studio_artifact:primary",
            output_payload={"content": "# Previous"},
        ),
        "studio_artifact:other": _FakeArtifact(
            id="studio_artifact:other",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Other revision",
            revision_of_id="studio_artifact:other-primary",
        ),
    }

    response = _client().get("/api/studio/artifacts/studio_artifact:primary/revisions")

    assert response.status_code == 200
    body = response.json()
    assert [item["id"] for item in body] == ["studio_artifact:revision"]
    assert body[0]["revision_of_id"] == "studio_artifact:primary"


def test_create_workflow_run_for_artifact_requires_approval(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
            source_ids=["source:one"],
        ),
    }

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:1/workflow-runs",
        json={
            "title": "Generate Report",
            "source_ids": ["source:one"],
            "approval_required": True,
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["id"] == "studio_workflow_run:1"
    assert body["artifact_id"] == "studio_artifact:1"
    assert body["notebook_id"] == "notebook:alpha"
    assert body["status"] == "awaiting_approval"
    assert body["approval_required"] is True
    assert body["steps"][0]["id"] == "context"
    assert body["steps"][1]["id"] == "privacy_gate"
    assert body["steps"][1]["status"] == "pending"


def test_create_workflow_run_without_approval_submits_generation_command(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:course-pack": _FakeArtifact(
            id="studio_artifact:course-pack",
            notebook_id="notebook:training",
            artifact_type="course_pack",
            title="Course Pack",
            source_ids=["source:ready"],
        ),
    }

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:ready"
            return SimpleNamespace(
                id=source_id,
                title="Ready transcript",
                full_text="Transcript text is ready.",
                command=None,
            )

    submitted: list[tuple[str, str, dict]] = []

    class _CommandServiceMock:
        @staticmethod
        async def submit_command_job(app_name, command_name, command_args):
            submitted.append((app_name, command_name, command_args))
            return "command:studio-generate"

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod, "CommandService", _CommandServiceMock, raising=False
    )
    monkeypatch.setattr(studio_mod, "provision_langchain_model", AsyncMock())

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:course-pack/workflow-runs",
        json={
            "title": "Generate Course Pack",
            "source_ids": ["source:ready"],
            "approval_required": False,
        },
    )

    assert response.status_code == 201
    body = response.json()
    assert body["status"] == "queued"
    assert body["command_id"] == "command:studio-generate"
    assert submitted == [
        (
            "open_notebook",
            "generate_studio_artifact",
            {
                "artifact_id": "studio_artifact:course-pack",
                "workflow_run_id": "studio_workflow_run:1",
            },
        )
    ]
    assert not studio_mod.provision_langchain_model.called


def test_create_workflow_run_rejects_not_ready_sources_before_queueing(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:course-pack": _FakeArtifact(
            id="studio_artifact:course-pack",
            notebook_id="notebook:training",
            artifact_type="course_pack",
            title="Course Pack",
            source_ids=["source:queued"],
        ),
    }

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:queued"
            return SimpleNamespace(
                id=source_id,
                title="Queued video",
                full_text=None,
                command="command:process-video",
            )

    class _CommandServiceMock:
        submit_command_job = AsyncMock(return_value="command:should-not-submit")

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod, "CommandService", _CommandServiceMock, raising=False
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:course-pack/workflow-runs",
        json={
            "title": "Generate Course Pack",
            "source_ids": ["source:queued"],
            "approval_required": False,
        },
    )

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "sources_not_ready"
    assert detail["not_ready_sources"] == [
        {
            "source_id": "source:queued",
            "title": "Queued video",
            "command_id": "command:process-video",
        }
    ]
    assert not _CommandServiceMock.submit_command_job.called
    assert _FakeWorkflowRun.records == {}


def test_list_workflow_runs_for_artifact(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
        ),
    }
    _FakeWorkflowRun.records = {
        "studio_workflow_run:1": _FakeWorkflowRun(
            id="studio_workflow_run:1",
            artifact_id="studio_artifact:1",
            notebook_id="notebook:alpha",
            title="Generate Report",
            status="queued",
        ),
    }

    response = _client().get("/api/studio/artifacts/studio_artifact:1/workflow-runs")

    assert response.status_code == 200
    assert [item["id"] for item in response.json()] == ["studio_workflow_run:1"]


def test_approve_workflow_run_releases_privacy_gate(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
            source_ids=["source:one"],
        ),
    }
    _FakeWorkflowRun.records = {
        "studio_workflow_run:1": _FakeWorkflowRun(
            id="studio_workflow_run:1",
            artifact_id="studio_artifact:1",
            notebook_id="notebook:alpha",
            title="Generate Report",
            status="awaiting_approval",
            approval_required=True,
            source_ids=["source:one"],
            steps=[
                {"id": "context", "label": "Context built", "status": "completed"},
                {"id": "privacy_gate", "label": "Privacy gate", "status": "pending"},
            ],
        ),
    }

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:one"
            return SimpleNamespace(
                id=source_id,
                title="Ready source",
                full_text="Ready source text.",
                command=None,
            )

    class _CommandServiceMock:
        @staticmethod
        async def submit_command_job(_app_name, _command_name, _command_args):
            return "command:studio-report"

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod, "CommandService", _CommandServiceMock, raising=False
    )

    response = _client().post("/api/studio/workflow-runs/studio_workflow_run:1/approve")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "queued"
    assert body["approval_required"] is False
    assert body["steps"][1]["status"] == "completed"


def test_approve_workflow_run_submits_generation_command(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _install_fake_workflow_runs(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:course-pack": _FakeArtifact(
            id="studio_artifact:course-pack",
            notebook_id="notebook:training",
            artifact_type="course_pack",
            title="Course Pack",
            source_ids=["source:ready"],
        ),
    }
    _FakeWorkflowRun.records = {
        "studio_workflow_run:1": _FakeWorkflowRun(
            id="studio_workflow_run:1",
            artifact_id="studio_artifact:course-pack",
            notebook_id="notebook:training",
            title="Generate Course Pack",
            status="awaiting_approval",
            approval_required=True,
            source_ids=["source:ready"],
            steps=[
                {"id": "context", "label": "Context built", "status": "completed"},
                {"id": "privacy_gate", "label": "Privacy gate", "status": "pending"},
                {"id": "model_route", "label": "Model route", "status": "blocked"},
                {
                    "id": "artifact_generation",
                    "label": "Course Pack",
                    "status": "blocked",
                },
            ],
        ),
    }

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:ready"
            return SimpleNamespace(
                id=source_id,
                title="Ready transcript",
                full_text="Transcript text is ready.",
                command=None,
            )

    submitted: list[tuple[str, str, dict]] = []

    class _CommandServiceMock:
        @staticmethod
        async def submit_command_job(app_name, command_name, command_args):
            submitted.append((app_name, command_name, command_args))
            return "command:studio-approved"

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod, "CommandService", _CommandServiceMock, raising=False
    )
    monkeypatch.setattr(studio_mod, "provision_langchain_model", AsyncMock())

    response = _client().post("/api/studio/workflow-runs/studio_workflow_run:1/approve")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "queued"
    assert body["approval_required"] is False
    assert body["command_id"] == "command:studio-approved"
    assert submitted == [
        (
            "open_notebook",
            "generate_studio_artifact",
            {
                "artifact_id": "studio_artifact:course-pack",
                "workflow_run_id": "studio_workflow_run:1",
            },
        )
    ]
    assert not studio_mod.provision_langchain_model.called


def test_update_artifact_patches_only_supplied_fields(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Draft",
            status="pending",
        )
    }

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:1",
        json={"title": "Final", "status": "completed"},
    )

    assert response.status_code == 200
    body = response.json()
    assert body["title"] == "Final"
    assert body["status"] == "completed"
    assert body["artifact_type"] == "report"


def test_update_artifact_rejects_invalid_structured_document(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:structured",
        notebook_id="notebook:alpha",
        artifact_type="quiz",
        title="Quiz",
        output_payload={"content": "# Existing"},
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:structured",
        json={
            "output_payload": {
                "schema_version": 1,
                "document": {"artifact_type": "quiz", "title": "Broken"},
                "markdown": "# Broken",
                "content": "# Broken",
            }
        },
    )

    assert response.status_code == 422
    assert response.json()["detail"]["code"] == "invalid_artifact_document"
    assert artifact.output_payload == {"content": "# Existing"}


def test_update_artifact_accepts_legacy_markdown_payload(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:legacy",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:legacy",
        json={"output_payload": {"content": "# Owner edited legacy artifact"}},
    )

    assert response.status_code == 200
    assert response.json()["output_payload"] == {
        "content": "# Owner edited legacy artifact"
    }


def test_update_artifact_rejects_unknown_schema_version(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:future",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:future",
        json={
            "output_payload": {
                "schema_version": 2,
                "document": {"artifact_type": "report", "title": "Future"},
                "content": "# Compatibility fallback",
            }
        },
    )

    assert response.status_code == 422
    assert response.json()["detail"]["code"] == "unsupported_artifact_schema"


def test_update_artifact_renders_canonical_markdown_and_keeps_extras(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:cards",
        notebook_id="notebook:alpha",
        artifact_type="flashcards",
        title="Cards",
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:cards",
        json={
            "output_payload": {
                "schema_version": 1,
                "document": {
                    "schema_version": 1,
                    "artifact_type": "flashcards",
                    "title": "Edited cards",
                    "cards": [{"front": "Question", "back": "Answer"}],
                },
                "markdown": "# Client supplied stale Markdown",
                "content": "# Different stale alias",
                "validation": {"status": "valid", "errors": []},
                "study_progress": {"index": 1},
            }
        },
    )

    assert response.status_code == 200
    output = response.json()["output_payload"]
    assert output["markdown"] == output["content"]
    assert output["markdown"].startswith("# Edited cards\n")
    assert "## Flashcard 1" in output["markdown"]
    assert "stale" not in output["markdown"]
    assert output["study_progress"] == {"index": 1}


def test_update_artifact_recomputes_derived_metadata(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:table",
        notebook_id="notebook:alpha",
        artifact_type="data_table",
        title="Table",
        citations=[{"marker": "[S1]", "source_id": "source:one"}],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:table",
        json={
            "output_payload": {
                "schema_version": 1,
                "document": {
                    "schema_version": 1,
                    "artifact_type": "data_table",
                    "title": "Edited table",
                    "columns": ["Topic", "Notes"],
                    "rows": [
                        {
                            "values": {"Topic": "New", "Notes": "Fresh"},
                            "citations": ["[S1]"],
                        }
                    ],
                },
                "markdown": "# Stale",
                "content": "# Stale",
                "validation": {"status": "valid", "errors": []},
                "data_table_rows": [{"Topic": "Old", "Notes": "Stale"}],
                "citation_warnings": {"unsupported_markers": ["[S9]"]},
                "study_progress": {"selected_row": 0},
            }
        },
    )

    assert response.status_code == 200
    output = response.json()["output_payload"]
    assert output["data_table_rows"] == [
        {"Topic": "New", "Notes": "Fresh", "Source": "[S1]"}
    ]
    assert "citation_warnings" not in output
    assert output["study_progress"] == {"selected_row": 0}


def test_update_visual_artifact_snapshots_and_refreshes_exports(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    previous_document = {
        "schema_version": 1,
        "artifact_type": "slide_deck",
        "title": "Original slides",
        "audience": "Researchers",
        "slides": [{"title": "Original slide", "bullets": ["Old content"]}],
    }
    artifact = _FakeArtifact(
        id="studio_artifact:edited-slides",
        notebook_id="notebook:alpha",
        artifact_type="slide_deck",
        title="Slides",
        status="completed",
        output_format="markdown",
        output_payload={
            "schema_version": 1,
            "document": previous_document,
            "markdown": "# Original slides\n",
            "content": "# Original slides\n",
            "validation": {"status": "valid", "errors": []},
        },
        export_paths={"pptx": str(tmp_path / "original.pptx")},
    )
    _FakeArtifact.records = {artifact.id: artifact}

    response = _client().patch(
        "/api/studio/artifacts/studio_artifact:edited-slides",
        json={
            "output_payload": {
                "schema_version": 1,
                "document": {
                    "schema_version": 1,
                    "artifact_type": "slide_deck",
                    "title": "Edited slides",
                    "audience": "Researchers",
                    "slides": [
                        {
                            "title": "Edited slide",
                            "bullets": ["Current content"],
                            "citations": ["[S1]"],
                        }
                    ],
                },
                "validation": {"status": "valid", "errors": []},
                "study_progress": {"selected_slide": 0},
            }
        },
    )

    assert response.status_code == 200
    body = response.json()
    assert body["output_payload"]["study_progress"] == {"selected_slide": 0}
    presentation = Presentation(body["export_paths"]["pptx"])
    slide_text = " ".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Edited slides" in slide_text
    assert "Edited slide" in slide_text
    revisions = [
        row
        for row in _FakeArtifact.records.values()
        if row.revision_of_id == artifact.id
    ]
    assert len(revisions) == 1
    assert revisions[0].output_payload["document"] == previous_document
    assert revisions[0].export_paths == {"pptx": str(tmp_path / "original.pptx")}


def test_delete_artifact_deletes_record(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    _FakeArtifact.records = {
        "studio_artifact:1": _FakeArtifact(
            id="studio_artifact:1",
            notebook_id="notebook:alpha",
            artifact_type="report",
            title="Report",
        )
    }

    response = _client().delete("/api/studio/artifacts/studio_artifact:1")

    assert response.status_code == 200
    assert response.json() == {"deleted": True, "id": "studio_artifact:1"}
    assert "studio_artifact:1" in _FakeArtifact.deleted


def test_generate_artifact_is_hidden_when_feature_flag_disabled(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "0")

    response = _client().post("/api/studio/artifacts/studio_artifact:1/generate")

    assert response.status_code == 404
    assert response.json()["detail"] == "Evidence Studio is not enabled"


def test_generate_artifact_uses_selected_sources_and_saves_markdown(
    monkeypatch, tmp_path
):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:1",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        source_ids=["source:one"],
        prompt="Make an executive report.",
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:one"
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    fake_chain = _json_chain("report", title="Generated Report")
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:1/generate")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "completed"
    assert body["output_format"] == "markdown"
    assert body["output_payload"]["content"].startswith("# Generated Report")
    assert body["output_payload"]["schema_version"] == 1
    assert body["output_payload"]["document"]["artifact_type"] == "report"
    assert body["output_payload"]["markdown"] == body["output_payload"]["content"]
    assert body["output_payload"]["validation"]["status"] == "valid"
    assert body["export_paths"]["markdown"].endswith(".md")
    assert body["export_paths"]["json"].endswith(".json")
    markdown_export = tmp_path / body["export_paths"]["markdown"].split("/")[-1]
    json_export = tmp_path / body["export_paths"]["json"].split("/")[-1]
    assert markdown_export.exists()
    assert json_export.exists()
    assert "# Generated Report" in markdown_export.read_text()
    assert "Source One" in markdown_export.read_text()
    exported_metadata = json.loads(json_export.read_text())
    assert exported_metadata["export_paths"] == body["export_paths"]
    assert exported_metadata["citations"][0]["source_id"] == "source:one"
    assert exported_metadata["citations"][0]["marker"] == "[S1]"
    assert body["citations"] == [
        {
            "source_id": "source:one",
            "title": "Source One",
            "marker": "[S1]",
            "location": "Source [S1]",
            "preview": "Important source text.",
        }
    ]
    assert _FakeArtifact.records["studio_artifact:1"].status == "completed"
    assert fake_chain.ainvoke.await_count == 1
    messages = fake_chain.ainvoke.call_args.args[0]
    assert "executive report" in messages[0].content.lower()
    assert "Use source markers like [S1]" in messages[0].content
    assert "## Source [S1]: Source One" in messages[1].content
    assert "Source ID: source:one" in messages[1].content
    assert "Important source text" in messages[1].content


def test_generate_artifact_rejects_podcast_audio_text_generation(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:podcast-audio",
        notebook_id="notebook:alpha",
        artifact_type="podcast_audio",
        title="Audio overview",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    provision = AsyncMock()
    monkeypatch.setattr(studio_mod, "provision_langchain_model", provision)

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:podcast-audio/generate"
    )

    assert response.status_code == 422
    assert "does not support structured generation" in response.json()["detail"]
    assert provision.await_count == 0


def test_generate_artifact_returns_409_when_sources_not_ready(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:not-ready",
        notebook_id="notebook:alpha",
        artifact_type="course_pack",
        title="Course Pack",
        source_ids=["source:queued"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:queued"
            return SimpleNamespace(
                id="source:queued",
                title="Queued source",
                full_text=None,
                command="command:process-source",
            )

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:not-ready/generate"
    )

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "sources_not_ready"
    assert detail["not_ready_sources"] == [
        {
            "source_id": "source:queued",
            "title": "Queued source",
            "command_id": "command:process-source",
        }
    ]
    assert _FakeArtifact.records["studio_artifact:not-ready"].status == "pending"
    assert not studio_mod.provision_langchain_model.called


def test_generate_course_pack_blocks_completed_source_without_text(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:empty-text",
        notebook_id="notebook:alpha",
        artifact_type="course_pack",
        title="Course Pack",
        source_ids=["source:empty"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            assert source_id == "source:empty"
            return SimpleNamespace(
                id="source:empty",
                title="Completed Empty Source",
                status="completed",
                full_text="   ",
                command=None,
            )

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:empty-text/generate"
    )

    assert response.status_code == 409
    detail = response.json()["detail"]
    assert detail["code"] == "sources_not_ready"
    assert detail["not_ready_sources"] == [
        {
            "source_id": "source:empty",
            "title": "Completed Empty Source",
            "command_id": None,
        }
    ]
    assert _FakeArtifact.records["studio_artifact:empty-text"].status == "pending"
    assert not studio_mod.provision_langchain_model.called


def test_generate_course_pack_saves_training_sidecar_exports(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:course-pack",
        notebook_id="notebook:training",
        artifact_type="course_pack",
        title="Admin Course Pack",
        source_ids=["source:video", "source:pdf"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            return SimpleNamespace(
                id=source_id,
                title="Training Source",
                full_text="A transcript and reference reading for admins.",
            )

    fake_chain = _json_chain(
        "course_pack",
        document={
            "artifact_type": "course_pack",
            "title": "Course Pack",
            "audience": "Workspace admins",
            "learning_outcomes": ["Map model roles to training tasks"],
            "modules": [
                {
                    "title": "Local Model Orientation",
                    "summary": "Duration: 20 minutes",
                    "lessons": [
                        {
                            "title": "Learner handout",
                            "content": "Compare source synthesis and study-fast roles.",
                            "duration_minutes": 20,
                            "exercise": "Select a model for course-pack generation.",
                            "facilitator_notes": (
                                "Open the local model settings page during the demo."
                            ),
                            "citations": ["[S1]"],
                        }
                    ],
                }
            ],
            "final_assessment": [
                {
                    "prompt": "Build a source-grounded learner handout.",
                    "options": [
                        {"id": "A", "text": "Use cited evidence"},
                        {"id": "B", "text": "Invent examples"},
                    ],
                    "correct_option_id": "A",
                    "citations": ["[S2]"],
                }
            ],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:course-pack/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["artifact_type"] == "course_pack"
    assert body["export_paths"]["instructor_guide"].endswith("-instructor-guide.md")
    assert body["export_paths"]["learner_handout"].endswith("-learner-handout.md")
    assert body["export_paths"]["module_checklist"].endswith("-module-checklist.json")
    assert body["export_paths"]["assessment"].endswith("-assessment.md")
    assert body["export_paths"]["scorm_package"].endswith("-scorm.zip")
    assert body["export_paths"]["xapi_package"].endswith("-xapi.zip")
    assert body["output_payload"]["course_pack_modules"] == [
        {
            "title": "Module 1: Local Model Orientation",
            "summary": "Duration: 20 minutes",
            "has_facilitator_notes": True,
        }
    ]

    instructor_export = (
        tmp_path / body["export_paths"]["instructor_guide"].split("/")[-1]
    )
    learner_export = tmp_path / body["export_paths"]["learner_handout"].split("/")[-1]
    checklist_export = (
        tmp_path / body["export_paths"]["module_checklist"].split("/")[-1]
    )
    assessment_export = tmp_path / body["export_paths"]["assessment"].split("/")[-1]
    assert "Open the local model settings page" in instructor_export.read_text()
    assert "Open the local model settings page" not in learner_export.read_text()
    checklist = json.loads(checklist_export.read_text())
    assert checklist["modules"][0]["title"] == "Module 1: Local Model Orientation"
    assert checklist["modules"][0]["has_facilitator_notes"] is True
    assert "Build a source-grounded learner handout" in assessment_export.read_text()

    scorm_export = tmp_path / body["export_paths"]["scorm_package"].split("/")[-1]
    xapi_export = tmp_path / body["export_paths"]["xapi_package"].split("/")[-1]
    with zipfile.ZipFile(scorm_export) as package:
        assert set(package.namelist()) >= {
            "imsmanifest.xml",
            "index.html",
            "instructor-guide.md",
            "learner-handout.md",
            "module-checklist.json",
            "assessment.md",
        }
        assert "Admin Course Pack" in package.read("imsmanifest.xml").decode()
        assert "Course Pack" in package.read("index.html").decode()
    with zipfile.ZipFile(xapi_export) as package:
        assert set(package.namelist()) >= {
            "tincan.xml",
            "xapi-statements.json",
            "instructor-guide.md",
            "learner-handout.md",
            "module-checklist.json",
            "assessment.md",
        }
        statements = json.loads(package.read("xapi-statements.json").decode())
        assert statements["activity"]["id"].endswith("studio_artifact:course-pack")
        assert statements["modules"][0]["title"] == "Module 1: Local Model Orientation"


def test_generate_course_pack_flags_unsupported_citation_markers(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:citation-guard",
        notebook_id="notebook:training",
        artifact_type="course_pack",
        title="Citation Guard Course Pack",
        source_ids=["source:one", "source:two"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, source_id):
            return SimpleNamespace(
                id=source_id,
                title=f"Source {source_id.rsplit(':', 1)[-1]}",
                full_text=f"Grounded content from {source_id}.",
            )

    fake_chain = _json_chain(
        "course_pack",
        document={
            "artifact_type": "course_pack",
            "title": "Course Pack",
            "audience": "Learners",
            "learning_outcomes": ["Verify citations"],
            "modules": [
                {
                    "title": "Verified markers",
                    "lessons": [
                        {
                            "title": "Citation check",
                            "content": "Use the first source for supported claims.",
                            "citations": ["[S1]", "[S3]"],
                        }
                    ],
                }
            ],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:citation-guard/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert [citation["source_id"] for citation in body["citations"]] == [
        "source:one",
        "source:two",
    ]
    assert [citation["marker"] for citation in body["citations"]] == ["[S1]", "[S2]"]
    assert body["output_payload"]["citation_warnings"] == {
        "unsupported_markers": ["[S3]"],
    }
    exported_metadata = json.loads(
        (tmp_path / body["export_paths"]["json"].split("/")[-1]).read_text()
    )
    assert exported_metadata["output_payload"]["citation_warnings"] == {
        "unsupported_markers": ["[S3]"],
    }


def test_generate_artifact_uses_role_routed_registered_model(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_MODEL_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:1",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}
    model = tmp_path / "GGUF" / "Qwen3-Coder-30B-A3B-Q4_K_M.gguf"
    model.parent.mkdir()
    model.write_bytes(b"x" * 4096)

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    class _ModelMock:
        @classmethod
        async def get_models_by_type(cls, model_type):
            assert model_type == "language"
            return [
                SimpleNamespace(
                    id="model:qwen-coder",
                    name="Qwen3-Coder-30B-A3B-Q4_K_M",
                    provider="openai_compatible",
                )
            ]

    captured: list[dict] = []
    fake_chain = _json_chain("report", title="Report")

    async def _fake_provision(content, model_id, default_type, **kwargs):
        captured.append(
            {
                "content": content,
                "model_id": model_id,
                "default_type": default_type,
                "kwargs": kwargs,
            }
        )
        return fake_chain

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(studio_mod, "Model", _ModelMock)
    monkeypatch.setattr(studio_mod, "provision_langchain_model", _fake_provision)

    response = _client().post("/api/studio/artifacts/studio_artifact:1/generate")

    assert response.status_code == 200
    body = response.json()
    assert body["model_id"] == "model:qwen-coder"
    assert body["provider"] == "openai_compatible"
    assert captured[0]["model_id"] == "model:qwen-coder"
    assert captured[0]["default_type"] == "chat"


def test_generate_artifact_keeps_explicit_model_over_role_routing(
    monkeypatch, tmp_path
):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_MODEL_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:manual",
        notebook_id="notebook:alpha",
        artifact_type="quiz",
        title="Quiz",
        source_ids=["source:one"],
        model_id="model:manual",
        provider="anthropic",
    )
    _FakeArtifact.records = {artifact.id: artifact}
    model = tmp_path / "GGUF" / "gemma-3-4b-it-Q4_K_M.gguf"
    model.parent.mkdir()
    model.write_bytes(b"x" * 4096)

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    class _ModelMock:
        @classmethod
        async def get_models_by_type(cls, _model_type):
            return [
                SimpleNamespace(
                    id="model:gemma",
                    name="gemma-3-4b-it-Q4_K_M",
                    provider="openai_compatible",
                )
            ]

    captured: list[str | None] = []
    fake_chain = _json_chain("quiz", title="Quiz")

    async def _fake_provision(_content, model_id, _default_type, **_kwargs):
        captured.append(model_id)
        return fake_chain

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(studio_mod, "Model", _ModelMock)
    monkeypatch.setattr(studio_mod, "provision_langchain_model", _fake_provision)

    response = _client().post("/api/studio/artifacts/studio_artifact:manual/generate")

    assert response.status_code == 200
    body = response.json()
    assert body["model_id"] == "model:manual"
    assert body["provider"] == "anthropic"
    assert captured == ["model:manual"]


def test_generate_artifact_preserves_previous_output_as_revision(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    previous_payload = {
        "schema_version": 1,
        "document": {
            "schema_version": 1,
            "artifact_type": "report",
            "title": "Previous Report",
            "sections": [{"heading": "Evidence", "body": "Previous answer."}],
        },
        "markdown": "# Previous Report\n",
        "content": "# Previous Report\n",
        "validation": {"status": "valid", "errors": []},
    }
    artifact = _FakeArtifact(
        id="studio_artifact:primary",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        status="completed",
        source_ids=["source:one"],
        prompt="Keep it concise.",
        model_id="model:local",
        provider="llamacpp",
        output_format="markdown",
        output_payload=previous_payload,
        citations=[
            {
                "source_id": "source:old",
                "title": "Old Source",
                "preview": "Old excerpt.",
            }
        ],
        export_paths={"markdown": "/tmp/report.md"},
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="New source text.",
            )

    fake_chain = _json_chain("report", title="New Report")
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:primary/generate")

    assert response.status_code == 200
    assert response.json()["output_payload"]["content"].startswith("# New Report")
    revisions = [
        row
        for row in _FakeArtifact.records.values()
        if row.revision_of_id == "studio_artifact:primary"
    ]
    assert len(revisions) == 1
    revision = revisions[0]
    assert revision.status == "completed"
    assert revision.title == "Report revision"
    assert revision.source_ids == ["source:one"]
    assert revision.prompt == "Keep it concise."
    assert revision.model_id == "model:local"
    assert revision.provider == "llamacpp"
    assert revision.output_format == "markdown"
    assert revision.output_payload == previous_payload
    assert revision.citations == [
        {
            "source_id": "source:old",
            "title": "Old Source",
            "preview": "Old excerpt.",
        }
    ]
    assert revision.export_paths == {"markdown": "/tmp/report.md"}


def test_generate_artifact_supports_first_text_artifact_types(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)

    for artifact_type, expected_instruction in [
        ("briefing", "short briefing"),
        ("faq", "source-grounded FAQ"),
        ("timeline", "chronological timeline"),
        ("course_pack", "instructor-ready Course Pack"),
        ("training_guide", "legacy alias for Course Pack"),
        ("flashcards", "flashcards"),
        ("quiz", "quiz"),
        ("data_table", "Data Table"),
    ]:
        artifact = _FakeArtifact(
            id=f"studio_artifact:{artifact_type}",
            notebook_id="notebook:alpha",
            artifact_type=artifact_type,
            title=artifact_type,
            source_ids=["source:one"],
        )
        _FakeArtifact.records = {artifact.id: artifact}

        fake_chain = _json_chain(artifact_type)
        monkeypatch.setattr(
            studio_mod,
            "provision_langchain_model",
            AsyncMock(return_value=fake_chain),
        )

        response = _client().post(f"/api/studio/artifacts/{artifact.id}/generate")

        assert response.status_code == 200
        assert response.json()["status"] == "completed"
        messages = fake_chain.ainvoke.call_args.args[0]
        assert expected_instruction in messages[0].content


def test_generate_artifact_supports_mind_map_instruction(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:mind-map",
        notebook_id="notebook:alpha",
        artifact_type="mind_map",
        title="Mind map",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="The project has source ingestion, citations, and local models.",
            )

    fake_chain = _json_chain("mind_map", title="Mind Map")
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:mind-map/generate")

    assert response.status_code == 200
    body = response.json()
    assert body["artifact_type"] == "mind_map"
    assert body["output_payload"]["content"].startswith("# Mind Map")
    messages = fake_chain.ainvoke.call_args.args[0]
    assert "mind map" in messages[0].content.lower()
    assert "nested markdown outline" in messages[0].content.lower()
    assert "relationships" in messages[0].content.lower()


def test_generate_artifact_supports_visual_study_artifact_instructions(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Evidence Studio creates research outputs with citations.",
            )

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)

    for artifact_type, expected_terms in [
        ("slide_deck", ("slide deck", "speaker notes", "slide")),
        ("infographic", ("infographic", "sections", "data callouts")),
    ]:
        artifact = _FakeArtifact(
            id=f"studio_artifact:{artifact_type}",
            notebook_id="notebook:alpha",
            artifact_type=artifact_type,
            title=artifact_type,
            source_ids=["source:one"],
        )
        _FakeArtifact.records = {artifact.id: artifact}

        fake_chain = _json_chain(artifact_type)
        monkeypatch.setattr(
            studio_mod,
            "provision_langchain_model",
            AsyncMock(return_value=fake_chain),
        )

        response = _client().post(f"/api/studio/artifacts/{artifact.id}/generate")

        assert response.status_code == 200
        assert response.json()["artifact_type"] == artifact_type
        messages = fake_chain.ainvoke.call_args.args[0]
        lower_prompt = messages[0].content.lower()
        for term in expected_terms:
            assert term in lower_prompt


def test_generate_artifact_supports_podcast_outline_instruction(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:podcast-outline",
        notebook_id="notebook:alpha",
        artifact_type="podcast_outline",
        title="Podcast outline",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Open Notebook Plus creates citation-backed research artifacts.",
            )

    fake_chain = _json_chain("podcast_outline", title="Podcast Outline")
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:podcast-outline/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["artifact_type"] == "podcast_outline"
    assert body["output_payload"]["content"].startswith("# Podcast Outline")
    messages = fake_chain.ainvoke.call_args.args[0]
    lower_prompt = messages[0].content.lower()
    assert "podcast outline" in lower_prompt
    assert "host segments" in lower_prompt
    assert "audio overview" in lower_prompt
    assert "citation markers" in lower_prompt


def test_generate_artifact_supports_research_run_instruction(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:research-run",
        notebook_id="notebook:alpha",
        artifact_type="research_run",
        title="Research run",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Open Notebook Plus needs competitive research synthesis.",
            )

    fake_chain = _json_chain("research_run", title="Research Run")
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:research-run/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["artifact_type"] == "research_run"
    assert body["output_payload"]["content"].startswith("# Research Run")
    messages = fake_chain.ainvoke.call_args.args[0]
    lower_prompt = messages[0].content.lower()
    assert "research run" in lower_prompt
    assert "multi-step" in lower_prompt
    assert "hypotheses" in lower_prompt
    assert "follow-up questions" in lower_prompt
    assert "citation markers" in lower_prompt


def test_generate_research_run_persists_structured_stage_metadata(
    monkeypatch, tmp_path
):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:research-stages",
        notebook_id="notebook:alpha",
        artifact_type="research_run",
        title="Research run",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Open Notebook Plus can route local models and generate artifacts.",
            )

    fake_chain = _json_chain(
        "research_run",
        document={
            "artifact_type": "research_run",
            "title": "Research Run",
            "objective": "Compare Open Notebook Plus with NotebookLM.",
            "hypotheses": ["Local model routing is a differentiator."],
            "stages": [
                {
                    "title": "Evidence-backed findings",
                    "findings": [
                        {
                            "text": "Evidence Studio can generate study artifacts.",
                            "citations": ["[S1]"],
                        }
                    ],
                }
            ],
            "next_actions": ["Which local model handles source synthesis best?"],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:research-stages/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["output_payload"]["document"]["objective"] == (
        "Compare Open Notebook Plus with NotebookLM."
    )
    stages = body["output_payload"]["research_stages"]
    assert stages == [
        {
            "title": "Hypotheses",
            "items": ["Local model routing is a differentiator."],
        },
        {
            "title": "Stage 1: Evidence-backed findings",
            "items": [
                "Status: complete",
                "Evidence Studio can generate study artifacts. [S1]",
            ],
        },
        {
            "title": "Next Actions",
            "items": ["Which local model handles source synthesis best?"],
        },
    ]
    json_export = tmp_path / body["export_paths"]["json"].split("/")[-1]
    exported_metadata = json.loads(json_export.read_text())
    assert exported_metadata["output_payload"]["research_stages"] == stages


def test_generate_data_table_persists_rows_and_csv_export(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:data-table",
        notebook_id="notebook:alpha",
        artifact_type="data_table",
        title="Data Table",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Open Notebook Plus supports local models and Evidence Studio.",
            )

    fake_chain = _json_chain(
        "data_table",
        document={
            "artifact_type": "data_table",
            "title": "Data Table",
            "columns": ["Topic", "Evidence", "Source", "Confidence", "Notes"],
            "rows": [
                {
                    "values": {
                        "Topic": "Local models",
                        "Evidence": "Scans AI_Models and routes roles [S1]",
                        "Source": "Source One",
                        "Confidence": "High",
                        "Notes": "User-owned runtime",
                    }
                },
                {
                    "values": {
                        "Topic": "Evidence Studio",
                        "Evidence": "Generates citation-backed artifacts [S1]",
                        "Source": "Source One",
                        "Confidence": "High",
                        "Notes": "Exportable",
                    }
                },
            ],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:data-table/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["artifact_type"] == "data_table"
    assert body["export_paths"]["csv"].endswith("-data-table.csv")
    assert body["output_payload"]["data_table_rows"] == [
        {
            "Topic": "Local models",
            "Evidence": "Scans AI_Models and routes roles [S1]",
            "Source": "Source One",
            "Confidence": "High",
            "Notes": "User-owned runtime",
        },
        {
            "Topic": "Evidence Studio",
            "Evidence": "Generates citation-backed artifacts [S1]",
            "Source": "Source One",
            "Confidence": "High",
            "Notes": "Exportable",
        },
    ]
    csv_export = tmp_path / body["export_paths"]["csv"].split("/")[-1]
    assert csv_export.exists()
    assert "Topic,Evidence,Source,Confidence,Notes" in csv_export.read_text()
    json_export = tmp_path / body["export_paths"]["json"].split("/")[-1]
    exported_metadata = json.loads(json_export.read_text())
    assert (
        exported_metadata["output_payload"]["data_table_rows"]
        == (body["output_payload"]["data_table_rows"])
    )
    messages = fake_chain.ainvoke.call_args.args[0]
    lower_prompt = messages[0].content.lower()
    assert "data table" in lower_prompt
    assert "markdown table" in lower_prompt
    assert "source marker" in lower_prompt


def test_generate_slide_deck_persists_pptx_and_pdf_exports(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:slides",
        notebook_id="notebook:alpha",
        artifact_type="slide_deck",
        title="Slides",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Source-grounded presentations retain evidence.",
            )

    fake_chain = _json_chain(
        "slide_deck",
        document={
            "artifact_type": "slide_deck",
            "title": "Evidence Slides",
            "audience": "Researchers",
            "slides": [
                {
                    "title": "Grounded output",
                    "bullets": ["Claims remain traceable."],
                    "speaker_notes": "Explain the retained evidence.",
                    "visual_direction": "Show a simple evidence path.",
                    "citations": ["[S1]"],
                }
            ],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:slides/generate")

    assert response.status_code == 200
    paths = response.json()["export_paths"]
    assert Path(paths["pptx"]).read_bytes().startswith(b"PK")
    assert Path(paths["pdf"]).read_bytes().startswith(b"%PDF")
    exported_metadata = json.loads(Path(paths["json"]).read_text())
    assert exported_metadata["export_paths"]["pptx"] == paths["pptx"]
    assert exported_metadata["export_paths"]["pdf"] == paths["pdf"]


def test_generate_infographic_persists_png_and_pdf_exports(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:infographic",
        notebook_id="notebook:alpha",
        artifact_type="infographic",
        title="Infographic",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Evidence can be summarized as a constrained visual.",
            )

    fake_chain = _json_chain(
        "infographic",
        document={
            "artifact_type": "infographic",
            "title": "Evidence Visual",
            "orientation": "portrait",
            "panels": [
                {
                    "kind": "metric",
                    "heading": "Coverage",
                    "value": "95%",
                    "body": "Resolved evidence",
                    "citations": ["[S1]"],
                }
            ],
        },
    )
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:infographic/generate"
    )

    assert response.status_code == 200
    paths = response.json()["export_paths"]
    assert Path(paths["png"]).read_bytes().startswith(b"\x89PNG")
    assert Path(paths["pdf"]).read_bytes().startswith(b"%PDF")


def test_visual_export_failure_keeps_completed_text_exports(monkeypatch, tmp_path):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    monkeypatch.setenv("DEEPER_NOTEBOOK_ARTIFACT_EXPORT_DIR", str(tmp_path))
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:visual-failure",
        notebook_id="notebook:alpha",
        artifact_type="slide_deck",
        title="Slides",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="PRIVATE SOURCE TEXT THAT MUST NOT ENTER WARNINGS",
            )

    def _broken_export(*_args, **_kwargs):
        raise RuntimeError("PRIVATE SOURCE TEXT THAT MUST NOT ENTER WARNINGS")

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=_json_chain("slide_deck")),
    )
    monkeypatch.setattr(
        studio_mod.artifact_generation_service,
        "export_slide_deck",
        _broken_export,
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:visual-failure/generate"
    )

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "completed"
    assert set(body["export_paths"]) == {"markdown", "json", "research_bundle"}
    warning = body["output_payload"]["export_warnings"]["visual"]
    assert warning["type"] == "RuntimeError"
    assert "PRIVATE SOURCE" not in json.dumps(warning)


def test_generate_artifact_uses_all_notebook_sources_when_none_selected(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:2",
        notebook_id="notebook:alpha",
        artifact_type="study_guide",
        title="Guide",
        source_ids=[],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _NotebookMock:
        @classmethod
        async def get(cls, notebook_id):
            assert notebook_id == "notebook:alpha"
            return cls()

        async def get_sources(self):
            return [
                SimpleNamespace(
                    id="source:a",
                    title="A",
                    full_text="Alpha text.",
                ),
                SimpleNamespace(
                    id="source:b",
                    title="B",
                    full_text="Beta text.",
                ),
            ]

    fake_chain = _json_chain("study_guide", title="Guide")
    monkeypatch.setattr(studio_mod, "Notebook", _NotebookMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:2/generate")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "completed"
    assert body["source_ids"] == ["source:a", "source:b"]
    assert "Alpha text" in fake_chain.ainvoke.call_args.args[0][1].content
    assert "Beta text" in fake_chain.ainvoke.call_args.args[0][1].content


def test_generate_artifact_returns_404_when_artifact_notebook_is_missing(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:missing-notebook",
        notebook_id="notebook:gone",
        artifact_type="study_guide",
        title="Guide",
        source_ids=[],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _NotebookMock:
        @classmethod
        async def get(cls, _notebook_id):
            raise NotFoundError("notebook not found")

    monkeypatch.setattr(studio_mod, "Notebook", _NotebookMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:missing-notebook/generate"
    )

    assert response.status_code == 404
    assert response.json()["detail"] == "Notebook not found: notebook:gone"
    assert _FakeArtifact.records["studio_artifact:missing-notebook"].status == "failed"
    assert not studio_mod.provision_langchain_model.called


def test_generate_artifact_marks_failed_when_model_errors(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:3",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    async def _boom(*_args, **_kwargs):
        raise RuntimeError("model exploded")

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(studio_mod, "provision_langchain_model", _boom)

    response = _client().post("/api/studio/artifacts/studio_artifact:3/generate")

    assert response.status_code == 502
    assert response.json()["detail"] == "Artifact generation failed"
    assert _FakeArtifact.records["studio_artifact:3"].status == "failed"


def test_generate_artifact_marks_failed_when_model_returns_blank_output(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:blank",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        source_ids=["source:one"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            return SimpleNamespace(
                id="source:one",
                title="Source One",
                full_text="Important source text.",
            )

    fake_chain = MagicMock()
    fake_chain.with_structured_output.side_effect = NotImplementedError
    fake_chain.ainvoke = AsyncMock(return_value=SimpleNamespace(content="   \n\t  "))
    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(return_value=fake_chain),
    )

    response = _client().post("/api/studio/artifacts/studio_artifact:blank/generate")

    assert response.status_code == 502
    assert response.json()["detail"] == "Artifact generation failed"
    saved = _FakeArtifact.records["studio_artifact:blank"]
    assert saved.status == "failed"
    assert saved.output_payload["validation"]["status"] == "invalid"
    assert "required structure" in saved.output_payload["error"].lower()
    assert saved.export_paths == {}


def test_generate_artifact_returns_404_when_selected_source_is_missing(monkeypatch):
    monkeypatch.setenv("DEEPER_NOTEBOOK_EVIDENCE_STUDIO", "1")
    _install_fake_artifacts(monkeypatch)
    artifact = _FakeArtifact(
        id="studio_artifact:missing-source",
        notebook_id="notebook:alpha",
        artifact_type="report",
        title="Report",
        source_ids=["source:gone"],
    )
    _FakeArtifact.records = {artifact.id: artifact}

    class _SourceMock:
        @classmethod
        async def get(cls, _source_id):
            raise NotFoundError("source not found")

    monkeypatch.setattr(studio_mod, "Source", _SourceMock)
    monkeypatch.setattr(
        studio_mod,
        "provision_langchain_model",
        AsyncMock(),
    )

    response = _client().post(
        "/api/studio/artifacts/studio_artifact:missing-source/generate"
    )

    assert response.status_code == 404
    assert response.json()["detail"] == "Source not found: source:gone"
    assert _FakeArtifact.records["studio_artifact:missing-source"].status == "failed"
    assert not studio_mod.provision_langchain_model.called
