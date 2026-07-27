"""Editable PPTX and deterministic PDF exports for slide documents."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from deeper_notebook.studio.exporters.theme import (
    BORDER,
    CORAL,
    INK,
    MUTED,
    NAVY,
    PAPER,
    TEAL,
    WHITE,
    draw_fitted_text,
    load_font,
    safe_pdf_title,
)
from deeper_notebook.studio.schemas import Slide, SlideDeckDocument

SLIDE_WIDTH = 1600
SLIDE_HEIGHT = 900


def _rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


def _set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)


def _add_title_slide(presentation: Presentation, document: SlideDeckDocument) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, NAVY)
    accent = slide.shapes.add_shape(
        1, Inches(0.7), Inches(0.85), Inches(0.12), Inches(5.8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(TEAL)
    accent.line.fill.background()
    _add_textbox(
        slide,
        left=1.15,
        top=1.4,
        width=10.8,
        height=2.5,
        text=document.title,
        size=34,
        color=WHITE,
        bold=True,
    )
    if document.audience:
        _add_textbox(
            slide,
            left=1.2,
            top=4.25,
            width=9.5,
            height=0.7,
            text=f"Prepared for {document.audience}",
            size=17,
            color="#D8E3ED",
        )
    _add_textbox(
        slide,
        left=1.2,
        top=6.65,
        width=6.0,
        height=0.4,
        text="OPEN NOTEBOOK PLUS  /  EVIDENCE STUDIO",
        size=10,
        color="#B7C8D8",
        bold=True,
    )
    slide.notes_slide.notes_text_frame.text = f"Deck title: {document.title}\nAudience: {document.audience or 'Not specified'}"


def _add_content_slide(
    presentation: Presentation,
    item: Slide,
    index: int,
    total: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, PAPER)
    rule = slide.shapes.add_shape(1, 0, 0, presentation.slide_width, Inches(0.12))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(TEAL if index % 2 else CORAL)
    rule.line.fill.background()
    _add_textbox(
        slide,
        left=0.72,
        top=0.48,
        width=11.7,
        height=0.8,
        text=item.title,
        size=27,
        color=INK,
        bold=True,
    )

    body = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.55), Inches(8.2), Inches(4.85)
    )
    frame = body.text_frame
    frame.clear()
    frame.word_wrap = True
    for bullet_index, bullet in enumerate(item.bullets):
        paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(20 if len(item.bullets) <= 5 else 17)
        paragraph.font.color.rgb = _rgb(INK)
        paragraph.space_after = Pt(14)

    if item.visual_direction:
        panel = slide.shapes.add_textbox(
            Inches(9.5), Inches(1.75), Inches(3.05), Inches(3.75)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = _rgb(WHITE)
        panel.line.color.rgb = _rgb(BORDER)
        panel.text_frame.clear()
        panel.text_frame.word_wrap = True
        heading = panel.text_frame.paragraphs[0]
        heading.text = "VISUAL DIRECTION"
        heading.font.name = "Arial"
        heading.font.bold = True
        heading.font.size = Pt(10)
        heading.font.color.rgb = _rgb(TEAL)
        detail = panel.text_frame.add_paragraph()
        detail.text = item.visual_direction
        detail.font.name = "Arial"
        detail.font.size = Pt(15)
        detail.font.color.rgb = _rgb(INK)
        detail.space_before = Pt(12)

    marker_text = " ".join(item.citations)
    if marker_text:
        _add_textbox(
            slide,
            left=0.85,
            top=6.78,
            width=8.0,
            height=0.35,
            text=f"Sources  {marker_text}",
            size=10,
            color=MUTED,
        )
    _add_textbox(
        slide,
        left=11.7,
        top=6.78,
        width=0.8,
        height=0.35,
        text=f"{index}/{total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    notes = []
    if item.speaker_notes:
        notes.append(f"Speaker notes:\n{item.speaker_notes}")
    if item.visual_direction:
        notes.append(f"Visual direction:\n{item.visual_direction}")
    if marker_text:
        notes.append(f"Sources: {marker_text}")
    slide.notes_slide.notes_text_frame.text = "\n\n".join(notes)


def _render_title_page(document: SlideDeckDocument) -> Image.Image:
    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), NAVY)
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((95, 110, 112, 745), radius=8, fill=TEAL)
    draw_fitted_text(
        draw,
        (155, 185),
        document.title,
        max_width=1260,
        max_height=260,
        max_size=76,
        min_size=46,
        fill=WHITE,
        bold=True,
        max_lines=3,
    )
    if document.audience:
        draw_fitted_text(
            draw,
            (160, 545),
            f"Prepared for {document.audience}",
            max_width=1100,
            max_height=80,
            max_size=30,
            min_size=22,
            fill="#D8E3ED",
        )
    draw.text(
        (160, 810),
        "OPEN NOTEBOOK PLUS  /  EVIDENCE STUDIO",
        font=load_font(18, bold=True),
        fill="#B7C8D8",
    )
    return image


def _render_content_page(item: Slide, index: int, total: int) -> Image.Image:
    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), PAPER)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, SLIDE_WIDTH, 14), fill=TEAL if index % 2 else CORAL)
    draw_fitted_text(
        draw,
        (85, 58),
        item.title,
        max_width=1400,
        max_height=100,
        max_size=45,
        min_size=30,
        fill=INK,
        bold=True,
        max_lines=2,
    )
    y = 205
    body_width = 930 if item.visual_direction else 1390
    for bullet in item.bullets[:8]:
        draw.ellipse((92, y + 10, 110, y + 28), fill=TEAL)
        _, height = draw_fitted_text(
            draw,
            (135, y),
            bullet,
            max_width=body_width,
            max_height=110,
            max_size=30,
            min_size=20,
            fill=INK,
            max_lines=3,
        )
        y += max(66, height + 28)
        if y > 720:
            break
    if item.visual_direction:
        draw.rounded_rectangle(
            (1090, 205, 1500, 650), radius=8, fill=WHITE, outline=BORDER, width=2
        )
        draw.text(
            (1130, 245),
            "VISUAL DIRECTION",
            font=load_font(18, bold=True),
            fill=TEAL,
        )
        draw_fitted_text(
            draw,
            (1130, 295),
            item.visual_direction,
            max_width=330,
            max_height=300,
            max_size=27,
            min_size=19,
            fill=INK,
            max_lines=8,
        )
    markers = " ".join(item.citations)
    if markers:
        draw.text(
            (90, 815),
            f"Sources  {markers}",
            font=load_font(18),
            fill=MUTED,
        )
    page = f"{index}/{total}"
    page_width = draw.textlength(page, font=load_font(18))
    draw.text((1510 - page_width, 815), page, font=load_font(18), fill=MUTED)
    return image


def render_slide_deck_images(
    document: SlideDeckDocument, output_dir: Path
) -> list[Path]:
    """Render deterministic 16:9 PNGs for private local media composition.

    The same renderer drives PDF export, so a Video Overview represents the
    reviewed slide artifact rather than a second, model-generated rendition.
    """
    if not isinstance(document, SlideDeckDocument):
        raise TypeError("render_slide_deck_images requires SlideDeckDocument")
    output_dir.mkdir(parents=True, exist_ok=True)
    pages = [_render_title_page(document)]
    pages.extend(
        _render_content_page(item, index, len(document.slides))
        for index, item in enumerate(document.slides, start=1)
    )
    paths: list[Path] = []
    try:
        for index, page in enumerate(pages, start=1):
            path = output_dir / f"slide-{index:03d}.png"
            page.save(path, "PNG", optimize=True)
            paths.append(path)
    finally:
        for page in pages:
            page.close()
    return paths


def export_slide_deck(
    document: SlideDeckDocument,
    pptx_path: Path,
    pdf_path: Path,
) -> None:
    if not isinstance(document, SlideDeckDocument):
        raise TypeError("export_slide_deck requires SlideDeckDocument")

    presentation = Presentation()
    presentation.slide_width = 12_192_000
    presentation.slide_height = 6_858_000
    _add_title_slide(presentation, document)
    for index, item in enumerate(document.slides, start=1):
        _add_content_slide(presentation, item, index, len(document.slides))
    presentation.core_properties.title = document.title
    presentation.core_properties.subject = "Evidence Studio slide deck"
    presentation.save(pptx_path)

    pages = [_render_title_page(document)]
    pages.extend(
        _render_content_page(item, index, len(document.slides))
        for index, item in enumerate(document.slides, start=1)
    )
    pages[0].save(
        pdf_path,
        "PDF",
        resolution=144.0,
        save_all=True,
        append_images=pages[1:],
        title=safe_pdf_title(document.title),
        author="Open Notebook Plus",
        subject="Evidence Studio slide deck",
    )
    for page in pages:
        page.close()


__all__ = ["export_slide_deck", "render_slide_deck_images"]
